from moviepy.editor import *
from moviepy.video.tools.drawing import color_gradient
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# === INPUTS ===
image_path = "/home/visharad/Downloads/liana-s-ird6OOE2LXI-unsplash.jpg"
audio_path = "/home/visharad/Downloads/videoplayback.wav"
slide_title = "Greenhouse Gases Explained"
slide_text = "Greenhouse gases trap heat in the atmosphere and keep Earth warm enough to sustain life. However, increased emissions from human activities like burning fossil fuels are enhancing this effect, leading to global warming."

# Monkey patch PIL.Image.ANTIALIAS
Image.ANTIALIAS = Image.LANCZOS

# === SETUP ===
audio = AudioFileClip(audio_path)
duration = audio.duration

# === IMAGE BACKGROUND ===
bg = ImageClip(image_path).set_duration(duration).resize(height=720).set_position("center")

# === TEXT OVERLAY ===
title_txt = TextClip(slide_title, fontsize=50, font='Liberation-Sans-Bold', color='white', bg_color='blue', size=(1280, None)).set_position(("center", 50)).set_duration(duration)
body_txt = TextClip(slide_text, fontsize=32, font='Liberation-Sans', color='white', size=(1000, None), method='caption').set_position(("center", 150)).set_duration(duration)

# === FINAL CLIP ===
final = CompositeVideoClip([bg, title_txt, body_txt]).set_audio(audio)
final.write_videofile("slide1.mp4", fps=24)

# === POWERPOINT GENERATION ===
# Create a presentation
prs = Presentation()

# Add a slide
slide_layout = prs.slide_layouts[5]  # Blank slide
slide = prs.slides.add_slide(slide_layout)

# Add title
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_tf = title_shape.text_frame
title_tf.text = slide_title  # Using the same title as video
title_tf.paragraphs[0].font.size = Pt(40)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add text
text_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
text_tf = text_shape.text_frame
text_tf.text = slide_text  # Using the same text as video
text_tf.paragraphs[0].font.size = Pt(24)

# Add image
slide.shapes.add_picture(image_path, Inches(1), Inches(5), width=Inches(8))

# Save
prs.save('slide_output.pptx')