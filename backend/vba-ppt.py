from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from colorthief import ColorThief
from PIL import Image
import requests # For downloading image from URL
import io # For handling image stream

# === CONFIG ===
title = "My Cool Slide"
text = "This slide was generated with Python on Linux!"
# Test with a URL
image_path = "https://plus.unsplash.com/premium_photo-1683749808307-e5597ac69f1e?w=500&auto=format&fit=crop&q=60&ixlib=rb-4.1.0&ixid=M3wxMjA3fDB8MHxzZWFyY2h8NXx8bWJhfGVufDB8fDB8fHww"
# Test with a local path (uncomment to test)
# image_path = "/path/to/your/local/image.jpg" 
output_path = "generated_slide.pptx"
template_path = "Modern shapes marketing plan presentation.pptx" # <--- ADD THIS: Path to your PowerPoint template

# === Utility ===

def get_image_stream_from_path(path_or_url):
    """
    Returns a file-like object (BytesIO stream) for an image,
    whether it's a local path or a URL.
    """
    if path_or_url.startswith(('http://', 'https://')):
        try:
            response = requests.get(path_or_url, stream=True)
            response.raise_for_status() # Raise an exception for bad status codes
            return io.BytesIO(response.content)
        except requests.exceptions.RequestException as e:
            print(f"Error downloading image from URL {path_or_url}: {e}")
            return None
    else:
        try:
            return open(path_or_url, 'rb') # Open local file in binary read mode
        except FileNotFoundError:
            print(f"Error: Local image file not found at {path_or_url}")
            return None
        except IOError as e:
            print(f"Error opening local image file {path_or_url}: {e}")
            return None


def get_dominant_color(image_input):
    """
    image_input can be a file path, URL, or a file-like object (stream).
    """
    if isinstance(image_input, str): # If it's a path or URL string
        stream = get_image_stream_from_path(image_input)
        if not stream:
            return (128, 128, 128) # Default grey if image fetch fails
    elif hasattr(image_input, 'read'): # If it's already a stream
        stream = image_input
        stream.seek(0) # Ensure stream is at the beginning
    else:
        print("Invalid image input for get_dominant_color")
        return (128, 128, 128)

    try:
        ct = ColorThief(stream)
        dominant_color = ct.get_color(quality=1)
        return dominant_color
    except Exception as e:
        print(f"Error getting dominant color: {e}")
        return (128, 128, 128) # Default grey on error
    finally:
        if isinstance(image_input, str) and stream and hasattr(stream, 'close'): # Close if we opened it
            stream.close()


def is_dark(rgb):
    r, g, b = rgb
    luminance = (0.299*r + 0.587*g + 0.114*b)/255
    return luminance < 0.5

def to_rgbcolor(rgb_tuple):
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])

def resize_and_center_image(prs, slide, image_input, max_width_in, max_height_in):
    """
    image_input can be a file path, URL, or a file-like object (stream).
    """
    # Slide size
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    if isinstance(image_input, str): # If it's a path or URL string
        stream = get_image_stream_from_path(image_input)
        if not stream:
            print("Could not add image to slide.")
            return
    elif hasattr(image_input, 'read'): # If it's already a stream
        stream = image_input
        stream.seek(0) # Ensure stream is at the beginning
    else:
        print("Invalid image input for resize_and_center_image")
        return

    try:
        # Open image using PIL from stream
        img = Image.open(stream)
        img_width_px, img_height_px = img.size
        dpi = img.info.get('dpi', (96, 96))[0] # Use image DPI if available, else default 96
        if dpi == 0: dpi = 96 # Handle cases where DPI might be 0

        img_width_in = img_width_px / dpi
        img_height_in = img_height_px / dpi

        # Scale proportionally
        width_ratio = max_width_in / img_width_in if img_width_in > 0 else 1
        height_ratio = max_height_in / img_height_in if img_height_in > 0 else 1
        scale = min(width_ratio, height_ratio, 1.0) # Don't scale up beyond 100% of max_width/height

        final_width_in = img_width_in * scale
        final_height_in = img_height_in * scale

        # Center it
        left_emu = (slide_width_emu - Inches(final_width_in)) / 2
        # Adjust top position slightly to account for title, or make it configurable
        top_offset_in = Inches(1.0) # Space for title and some margin
        top_emu = (slide_height_emu - Inches(final_height_in)) / 2 + top_offset_in

        # Ensure stream is at the beginning before passing to add_picture
        stream.seek(0)
        slide.shapes.add_picture(stream, left_emu, top_emu, width=Inches(final_width_in), height=Inches(final_height_in))
    
    except Exception as e:
        print(f"Error processing or adding image: {e}")
    finally:
        if isinstance(image_input, str) and stream and hasattr(stream, 'close'): # Close if we opened it
             stream.close()


# === MAIN ===
# Get a single image stream to pass around, avoiding multiple downloads/opens
image_stream = get_image_stream_from_path(image_path)

if image_stream:
    dominant_rgb = get_dominant_color(image_stream) # Pass the stream
    text_color = RGBColor(255, 255, 255) if is_dark(dominant_rgb) else RGBColor(0, 0, 0)
    bg_color = to_rgbcolor(dominant_rgb)

    # Load presentation from template if path is provided, otherwise create new
    if template_path:
        try:
            prs = Presentation(template_path)
            print(f"📄 Using template: {template_path}")
        except Exception as e:
            print(f"Error loading template '{template_path}': {e}. Creating a blank presentation instead.")
            prs = Presentation()
    else:
        prs = Presentation()
    
    # Use a blank layout (index 6 is usually blank, check your template's layouts)
    # You might want to choose a specific layout from your template by its index
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # Background
    # This will override the template's background for this specific slide.
    # Comment out if you want to keep the template's slide background.
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.add_paragraph()
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = text_color
    # p.alignment = PP_ALIGN.CENTER # Requires: from pptx.enum.text import PP_ALIGN

    # Body
    # Position body text below title, adjust as needed
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.2)) 
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    p = body_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18) # Adjusted for better fit
    p.font.color.rgb = text_color

    # Smart Image Insert - image will be placed considering title and body
    # Max width/height for image area, adjust top_offset in resize_and_center_image if needed
    resize_and_center_image(prs, slide, image_stream, max_width_in=8.0, max_height_in=4.0) 

    # Save
    prs.save(output_path)
    print(f"🧠 Slide created at: {output_path}")

    if hasattr(image_stream, 'close'):
        image_stream.close()
else:
    print(f"Could not process image from: {image_path}. Slide not generated.")

