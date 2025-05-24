from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from colorthief import ColorThief
from PIL import Image

# Load color from image
image_path = "/home/visharad/Downloads/liana-s-ird6OOE2LXI-unsplash.jpg"
color_thief = ColorThief(image_path)
dominant_color = color_thief.get_color(quality=1)

# Initialize presentation
prs = Presentation()

# Set slide width/height if needed
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Title slide layout (0 is usually title slide)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Auto-Generated Presentation"
subtitle.text = "With Custom Style"

# Add a colored rectangle as a design element
left = top = Inches(0)
width = prs.slide_width
height = Inches(1)

shape = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(*dominant_color)

# Save your template-like PPTX
prs.save("my_template.pptx")
